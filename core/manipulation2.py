
# pip install boto3 python-pptx

import boto3
import io
import os
import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
from pydub import AudioSegment
from concurrent.futures import ThreadPoolExecutor
import uuid

from ssml_validation import *

from pdf2image import convert_from_path

from pydub import AudioSegment
from pydub.playback import play


load_dotenv()

# Get environment variables
aws_access_key_id = os.getenv('aws_access_key_id')
aws_secret_access_key = os.getenv('aws_secret_access_key')
bucket_name = os.getenv('bucket_name')
region = os.getenv('region')


# Generate 0.5 seconds of silence
half_sec_silence = AudioSegment.silent(
    duration=500)  # duration in milliseconds


def split_input_path(input_path):
    usermail, filename = input_path.split('/', 1)
    return usermail, filename


def download_pptx_from_s3(usermail, project, filename):
    s3 = boto3.resource('s3', aws_access_key_id=aws_access_key_id,
                        aws_secret_access_key=aws_secret_access_key, region_name=region)
    obj = s3.Object(bucket_name, f'{usermail}/{project}/{filename}')
    pptx_file = io.BytesIO()
    obj.download_fileobj(pptx_file)
    return pptx_file


def upload_pptx_to_s3(usermail, project, filename, pptx_file):
    s3 = boto3.resource('s3', aws_access_key_id=aws_access_key_id,
                        aws_secret_access_key=aws_secret_access_key, region_name=region)
    # Include the 'edited' directory in the path
    obj = s3.Object(bucket_name, f'{usermail}/{project}/edited/{filename}')
    pptx_file.seek(0)
    obj.upload_fileobj(pptx_file)


def generate_tts(text, voice_id):
    polly_client = boto3.client('polly', aws_access_key_id=aws_access_key_id,
                                aws_secret_access_key=aws_secret_access_key, region_name=region)
    try:

        response = polly_client.synthesize_speech(
            VoiceId=voice_id, OutputFormat='mp3', Text=text, TextType='ssml', Engine='neural')
        # Create unique filename for the audio file
        filename = f'tts_{uuid.uuid4()}.mp3'
        with open(filename, 'wb') as out:  # open for [w]riting as [b]inary
            out.write(response['AudioStream'].read())
    except Exception as e:
        print(f"Error during audio generation: {str(e)}")
        return None

    return filename


def combine_audio_files(audio_files):
    combined = AudioSegment.empty()
    for audio_file in audio_files:
        audio = AudioSegment.from_mp3(audio_file)
        combined += audio
    combined_filename = f"/tmp/combined_{hash(''.join(audio_files))}.mp3"
    combined.export(combined_filename, format="mp3")
    return combined_filename


def process_slide(slide):
    try:
        notes_slide = slide.notes_slide
        # print(f"notes_slide: {notes_slide}")
        notes_text = notes_slide.notes_text_frame.text
        # print(f"notes_text: {notes_text}")
        if notes_text and notes_text.strip():  # Check if notes_text is not empty
            # Correct special characters and validate the SSML
            # print(f"not empty")
            checked_missing_tags = find_missing_tags(notes_text)
            # print(f"checked_missing_tags: {checked_missing_tags}")
            try:
                corrected_ssml = correct_special_characters(
                    checked_missing_tags)
            except Exception as e:
                print(
                    f"Error during SSML correct_special_characters: {str(e)}")
                return
            # print("Corrected SSML: \n", corrected_ssml)

            validation_result = validate_ssml(corrected_ssml, schema_path)
            # print("Validation result: ", validation_result)

            # # Test SSML validation function
            if not validation_result:
                raise ValueError
                # print(" NOT VALIDATE SSML")
            # Process the SSML
            # print("VALIDATE SSML")
            try:
                parsed_ssml = parse_ssml(corrected_ssml)
            except Exception as e:
                print(f"Error during SSML parsing: {str(e)}")
                return
            # print(f"parsed_ssml: {parsed_ssml}")
            # combined_text = ''.join(text for voice_name, text in parsed_ssml)
            # print(f"Combined text: {combined_text}")
            try:
                if len(parsed_ssml) == 1:
                    print("len == 1")
                    for voice_name, text in parsed_ssml:
                        filename = generate_tts(text, voice_name)
                        print(f"filename: {filename}")
                        audio = AudioSegment.from_file(filename, format='mp3')
                        print(f"audio: {audio}")
                        try:
                            # Add the combined audio to the slide
                            left, top, width, height = Inches(
                                1), Inches(2.5), Inches(1), Inches(1)
                            slide.shapes.add_movie(
                                filename, left, top, width, height, mime_type="audio/mp3", poster_frame_image=None)
                        except Exception as e:
                            print(
                                f"Error during audio combining/exporting or adding to slide: {str(e)}")
                            return

                        # Remove the temporary files after using them
                        try:
                            os.remove(filename)
                        except OSError as e:
                            print(f"Error: {e.filename} - {e.strerror}.")
                else:
                    audios = []
                    # Generate an mp3 file for each voice and text
                    for voice_name, text in parsed_ssml:
                        # generate_tts now outputs a file path
                        filename = generate_tts(text, voice_name)
                        print(f"filename: {filename}")
                        audio = AudioSegment.from_mp3(filename)
                        audios.append(audio)
                        audios.append(half_sec_silence)  # 0.5 seconds pause
                        # combined_text += text
                        # print(f"combined_text: {combined_text}")
                        # print(f"audios final: {audios}")
                    # Remove last silence segment
                    audios.pop()
                    # print(f"audios final: {audios}")
                    # Combine all audios into one
                    combined_audio = audios[0]
                    for audio in audios[1:]:
                        combined_audio += audio

                    # Export combined audio to a file
                    combined_filename = f'combined_{uuid.uuid4()}.mp3'
                    combined_audio.export(
                        combined_filename, format="mp3", bitrate="320k")

                    try:
                        # Add the combined audio to the slide
                        left, top, width, height = Inches(
                            1), Inches(2.5), Inches(1), Inches(1)
                        slide.shapes.add_movie(
                            combined_filename, left, top, width, height, mime_type="audio/mp3", poster_frame_image=None)
                    except Exception as e:
                        print(
                            f"Error during audio combining/exporting or adding to slide: {str(e)}")
                        return

                        # Remove the temporary files after using them
                    try:
                        os.remove(combined_filename)
                    except OSError as e:
                        print(f"Error: {e.filename} - {e.strerror}.")
            except Exception as e:
                print(f"Error during audio generation: {str(e)}")
                return  # Skip the slide if there was an error

    except Exception as e:
        print(f"Error during SSML correction/validation/parsing: {str(e)}")
        return


def add_tts_to_pptx(pptx_file):
    pptx_file.seek(0)
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        process_slide(slide)
    pptx_file.seek(0)
    prs.save(pptx_file)


def process_pptx(usermail, project, filename):
    print(usermail, project, filename)
    pptx_file = download_pptx_from_s3(usermail, project, filename)
    add_tts_to_pptx(pptx_file)
    edited_filename = f"{os.path.splitext(filename)[0]}_edited{os.path.splitext(filename)[1]}"
    upload_pptx_to_s3(usermail, project, edited_filename, pptx_file)


def process_preview(text):
    try:
        checked_missing_tags = find_missing_tags(text)
        corrected_ssml = correct_special_characters(checked_missing_tags)
        validation_result = validate_ssml(corrected_ssml, schema_path)
    except Exception as e:
        print(f"Error during SSML correction/validation: {str(e)}")
        return

    if validation_result:
        try:
            # continue processing if SSML validation passed
            parsed_ssml = parse_ssml(corrected_ssml)
        except Exception as e:
            print(f"Error during SSML parsing: {str(e)}")
            return

        if len(parsed_ssml) == 1:
            # print("len == 1")
            for voice_name, text in parsed_ssml:
                filename = generate_tts(text, voice_name)
                print(f"filename: {filename}")
                audio = AudioSegment.from_file(filename, format='mp3')
                print(f"audio: {audio}")
                try:
                    # Add the combined audio to the slide
                    left, top, width, height = Inches(
                        1), Inches(2.5), Inches(1), Inches(1)
                    slide.shapes.add_movie(
                        filename, left, top, width, height, mime_type="audio/mp3", poster_frame_image=None)
                except Exception as e:
                    print(
                        f"Error during audio combining/exporting or adding to slide: {str(e)}")
                    return

                # Remove the temporary files after using them
                try:
                    os.remove(filename)
                except OSError as e:
                    print(f"Error: {e.filename} - {e.strerror}.")
        else:
            audios = []
            # Generate an mp3 file for each voice and text
            for voice_name, text in parsed_ssml:
                # generate_tts now outputs a file path
                filename = generate_tts(text, voice_name)
                print(f"filename: {filename}")
                audio = AudioSegment.from_mp3(filename)
                audios.append(audio)
                audios.append(half_sec_silence)  # 0.5 seconds pause
                # combined_text += text
                # print(f"combined_text: {combined_text}")
                # print(f"audios final: {audios}")
                # Remove last silence segment
                audios.pop()
                # print(f"audios final: {audios}")
                # Combine all audios into one
                combined_audio = audios[0]
                for audio in audios[1:]:
                    combined_audio += audio

                    # Export combined audio to a file
                    combined_filename = f'combined_{uuid.uuid4()}.mp3'
                    combined_audio.export(
                        combined_filename, format="mp3", bitrate="320k")

                    try:
                        # Add the combined audio to the slide
                        left, top, width, height = Inches(
                            1), Inches(2.5), Inches(1), Inches(1)
                        slide.shapes.add_movie(
                            combined_filename, left, top, width, height, mime_type="audio/mp3", poster_frame_image=None)
                    except Exception as e:
                        print(
                            f"Error during audio combining/exporting or adding to slide: {str(e)}")
                        return

                        # Remove the temporary files after using them
                    try:
                        os.remove(combined_filename)
                    except OSError as e:
                        print(f"Error: {e.filename} - {e.strerror}.")
    else:
        print("SSML validation failed, cannot process the preview.")


def pptx_to_pdf(pptx_file):
    output_pdf_path = f"/tmp/{os.path.basename(pptx_file.name)}.pdf"
    os.system(
        f'libreoffice --headless --convert-to pdf {pptx_file.name} --outdir /tmp')
    return output_pdf_path


def pdf_to_images(pdf_path):
    return convert_from_path(pdf_path)


def extract_notes_from_slide(slide):
    slide_number = slide.slide_id
    notes_text = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
    return slide_number, notes_text


def extract_notes_and_images_from_pptx(pptx_file):
    # Save the pptx_file to a temporary file, so that it can be converted to pdf
    pptx_file_path = f"/tmp/{filename}"
    with open(pptx_file_path, 'wb') as f:
        f.write(pptx_file.read())

    pdf_file_path = pptx_to_pdf(pptx_file_path)
    images = pdf_to_images(pdf_file_path)

    pptx_file.seek(0)
    prs = Presentation(pptx_file)

    notes_and_images = []
    for i, slide in enumerate(prs.slides):
        slide_id, notes_text = extract_notes_from_slide(slide)
        tts_filename = process_preview(
            notes_text) if notes_text.strip() else None
        notes_and_images.append({
            "slide_id": slide_id,
            "notes": notes_text,
            "image": images[i],
            "tts": tts_filename
        })

    return notes_and_images


def process_pptx_split(usermail, project, filename):
    # Download the pptx file from S3
    pptx_file = download_pptx_from_s3(usermail, project, filename)

    # Extract the slide images, notes and tts for each slide
    slide_data = extract_notes_and_images_from_pptx(pptx_file)

    for i, data in enumerate(slide_data):
        slide_id, notes_text, image, tts = data["slide_id"], data["notes"], data["image"], data["tts"]

        # Save the image locally
        image_filename = f"/tmp/slide_{slide_id}.jpg"
        image.save(image_filename)

        # Save the tts locally
        tts_filename = f"/tmp/slide_{slide_id}_tts.mp3"
        with open(tts_filename, "wb") as f:
            f.write(tts)

        # Replace the image and tts in the data dictionary with their S3 URIs
        data["image"] = s3_image_uri
        data["tts"] = s3_tts_uri

    # Return the slide data (which now contains the URIs for the image and tts) as a JSON string
    return json.dumps(slide_data)
