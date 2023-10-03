import shutil
import base64
import os
import io
from pdf2image import convert_from_bytes, convert_from_path
from pydub import AudioSegment
import tempfile
from pptx.util import Inches
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx import Presentation
from ssml_validation import *
from exceptions import *


def zip_pptx(pptx_buffer):
    edited_buffer = BytesIO()
    with zipfile.ZipFile(edited_buffer, 'w') as zip_ref:
        zip_ref.writestr("file.pptx", pptx_buffer.getvalue())
    return edited_buffer


def unzip_file(zip_byte_data_io):
    with zipfile.ZipFile(zip_byte_data_io, 'r') as zip_ref:
        for name in zip_ref.namelist():
            if name.endswith('.pptx'):
                return BytesIO(zip_ref.read(name))
    return None


def create_folder(folder_path):
    os.makedirs(folder_path, exist_ok=True)
    return os.path.abspath(folder_path)


def delete_folder(folder):
    try:
        shutil.rmtree(folder)
    except Exception as e:
        raise ElaborationException(
            f"Critical: Exception while deleting temp folder: {str(e)}")


def combine_audios(audios):
    audios.pop()
    combined_audio = audios[0]
    for audio in audios[1:]:
        combined_audio += audio
    return combined_audio


def add_audio_to_slide(slide, audio_path):
    left, top, width, height = Inches(1), Inches(1.5), Inches(1), Inches(1)
    slide.shapes.add_movie(audio_path, left, top, width,
                           height, mime_type="audio/mp3")


def check_correct_validate_parse_text(notes_text):
    try:
        checked_missing_tags = find_missing_tags(notes_text)
        corrected_ssml = correct_special_characters(checked_missing_tags)
        validate_ssml(corrected_ssml)
        return parse_ssml(corrected_ssml)
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise ElaborationException(
            f"Exception during SSML correction/validation/parsing: {str(e)}")


def check_slide_have_notes(notes_slide):
    if notes_slide and notes_slide.notes_text_frame:
        notes_text = notes_slide.notes_text_frame.text
        if notes_text and notes_text.strip():
            return notes_text
    return None


def audiosegment_to_stream(audio_segment):
    try:
        stream = io.BytesIO()
        audio_segment.export(stream, format="mp3")
        stream.seek(0)
        return stream
    except Exception as e:
        raise ElaborationException(
            f"Exception while converting audio to stream: {str(e)}")

# END


def image_to_base64(image):
    try:
        image_bytes = io.BytesIO()
        image.save(image_bytes, format='JPEG')
        return base64.b64encode(image_bytes.getvalue()).decode('utf-8')
    except Exception as e:
        raise ElaborationException(
            f"Exception while converting image to base64: {str(e)}")


def split_input_path(input_path):
    usermail, filename = input_path.split('/', 1)
    return usermail, filename


def pdf_to_images(pdf_buffer):
    # Convert the PDF bytes to Pillow images
    images = convert_from_bytes(pdf_buffer.getvalue())

    # Convert each Pillow image to a BytesIO buffer containing the JPEG data
    buffers = []
    for image in images:
        buffer = io.BytesIO()
        image.save(buffer, format='JPEG')
        buffer.seek(0)  # Important! Reset the buffer's position to the start
        buffers.append(buffer)

    return buffers


def is_pptx_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    return file_extension.lower() == ".pptx"


def file_ispptx_exists_readpermission(file_path):
    return is_pptx_file and os.path.exists(file_path) and os.access(file_path, os.R_OK)


def slide_split_data(index, image_base64, audio_base64):
    return {
        "slide_id": index,
        "image": {
            "data": image_base64,
            "filename": f"slide_{index}.jpg"
        },
        "tts": {
            "data": audio_base64,
            "filename": f"slide_{index}.mp3"
        } if audio_base64 is not None else None
    }


def extract_image_from_slide(index, folder, image):
    image_path = os.path.join(folder, f'slide_{index}.jpg')
    image.save(image_path)
    return image_to_base64(image)
