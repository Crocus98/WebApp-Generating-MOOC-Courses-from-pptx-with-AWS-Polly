from pptx import Presentation
from pydub import AudioSegment
from dotenv import load_dotenv
from zipfile import ZipFile
from ssml_validation import *
from exceptions import *
from utils import *
import subprocess
import boto3
import uuid
import io
import os
import traceback
import zipfile

load_dotenv()


class S3Singleton:
    _instance = None

    def __new__(cls, aws_access_key_id, aws_secret_access_key, region_name, bucket_name):
        if cls._instance is None:
            cls._instance = super(S3Singleton, cls).__new__(cls)
            cls._instance.s3 = boto3.resource('s3',
                                              aws_access_key_id=aws_access_key_id,
                                              aws_secret_access_key=aws_secret_access_key,
                                              region_name=region_name)
            cls._instance.bucket_name = bucket_name
        return cls._instance


class Polly:
    def __init__(self, aws_access_key_id, aws_secret_access_key, region_name):
        self.polly = boto3.client('polly',
                                  aws_access_key_id=aws_access_key_id,
                                  aws_secret_access_key=aws_secret_access_key,
                                  region_name=region_name)

    def __del__(self):
        self.polly = None


aws_access_key_id = os.getenv('aws_access_key_id')
aws_secret_access_key = os.getenv('aws_secret_access_key')
bucket_name = os.getenv('bucket_name')
region = os.getenv('region')
schema_path = os.getenv('schema_path')
path_to_libreoffice = os.getenv("path_to_libreoffice")

s3_singleton = S3Singleton(
    aws_access_key_id, aws_secret_access_key, region, bucket_name)
polly_object = Polly(aws_access_key_id, aws_secret_access_key, region)


half_sec_silence = AudioSegment.silent(duration=500)


def zip_pptx(pptx_file, temp_folder, edited_filename):
    # edited_filename = f"{os.path.splitext(os.path.basename(pptx_file))[0]}_edited.zip"
    edited_filepath = os.path.join(temp_folder, edited_filename)

    with zipfile.ZipFile(edited_filepath, 'w') as zip_ref:
        zip_ref.write(pptx_file, os.path.basename(pptx_file))

    print(f"Uploading edited pptx: {edited_filename}")
    print(f"Creating zip at: {edited_filename}")

    return edited_filepath


def unzip_file_to_temp(zip_byte_data_io, usermail, project):
    # Determine the directory of the current script
    script_directory = os.path.dirname(os.path.abspath(__file__))

    # Create a temporary folder in the script directory
    temp_folder = os.path.join(script_directory, f"{usermail}_{project}_temp")
    if not os.path.exists(temp_folder):
        os.makedirs(temp_folder)

    # Save byte data to a temporary file
    temp_zip_file = os.path.join(temp_folder, "temp.zip")
    with open(temp_zip_file, 'wb') as f:
        f.write(zip_byte_data_io.getvalue())  # Use getvalue() to extract bytes

    # Unzip the saved temp file into the temp folder
    with zipfile.ZipFile(temp_zip_file, 'r') as zip_ref:
        zip_ref.extractall(temp_folder)

    # Assuming there's only one PPTX file inside the zip
    pptx_file = [os.path.join(temp_folder, f) for f in os.listdir(
        temp_folder) if f.endswith('.pptx')][0]

    return pptx_file, temp_folder


def download_file_from_s3(usermail, project, filename):
    try:
        obj = s3_singleton.s3.Object(
            bucket_name, f'{usermail}/{project}/{filename}')
        try:
            obj.get()
        except Exception as e:
            raise UserParameterException(
                f"File not found: check parameters. Error: {str(e)}")
        file = io.BytesIO()
        obj.download_fileobj(file)
        if file is None:
            raise AmazonException("file is None")
        return file
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise AmazonException(f"Exception downloading file from s3: {str(e)}")


def upload_file_to_s3(usermail, project, filename, file):
    try:
        obj = s3_singleton.s3.Object(
            bucket_name, f'{usermail}/{project}/edited/{filename}')
        #if (is_pptx_file(filename)):
        file.seek(0)
        obj.upload_fileobj(file)
        response = obj.get()
        if not (response and response['ResponseMetadata']['HTTPStatusCode'] == 200):
            raise AmazonException("Upload to S3 failed")
    except Exception as e:
        raise AmazonException(f"Exception uploading pptx to s3: {str(e)}")


def generate_tts(text, voice_id, filename):
    try:
        try:
            response = polly_object.polly.synthesize_speech(
                VoiceId=voice_id, OutputFormat='mp3', Text=text, TextType='ssml', Engine='neural')
            if not response['ResponseMetadata']['HTTPStatusCode'] == 200:
                raise AmazonException(
                    f"Polly failed to elaborate tts. Response is not 200.")
        except Exception as e:
            raise AmazonException(f"Exception from Polly: {str(e)}")
        with open(filename, 'wb') as out:
            out.write(response['AudioStream'].read())
    except AmazonException as e:
        raise AmazonException(e)
    except Exception as e:
        raise ElaborationException(
            f"Exception while saving audio into file: {str(e)}")


def pptx_to_pdf(pptx_file_path):
    if not file_ispptx_exists_readpermission(pptx_file_path):
        raise ElaborationException(
            f"Temporary copy of file is either not found, not a pptx or not readable due to permissions: {pptx_file_path}")

    print(f"Converting pptx to pdf: {pptx_file_path}")
    output_folder = os.path.dirname(pptx_file_path)
    output_pdf_path = os.path.join(
        output_folder, f"{os.path.splitext(os.path.basename(pptx_file_path))[0]}.pdf")
    print(f"Output pdf path: {output_pdf_path}")

    # Check if PDF already exists
    if os.path.exists(output_pdf_path):
        print(f"PDF already exists at {output_pdf_path}. Overwriting.")
        os.remove(output_pdf_path)

    #command = f"{path_to_libreoffice} --headless --convert-to pdf --outdir \"{output_folder}\" \"{pptx_file_path}\""
    command = f"/instdir/program/soffice --headless --invisible --nodefault --nofirststartwizard --nolockcheck --nologo --norestore --convert-to pdf --outdir \"{output_folder}\" \"{pptx_file_path}\""

    process = subprocess.Popen(
        command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, shell=True)

    stdout, stderr = process.communicate(timeout=60)
    if process.returncode != 0:
        print(f"STDOUT from LibreOffice: {stdout.decode()}")
        print(f"STDERR from LibreOffice: {stderr.decode()}")
        raise ElaborationException(
            f"Error converting PPTX to PDF using LibreOffice.")

    # Check if the output PDF exists
    if not os.path.exists(output_pdf_path):
        print("NOT OS!")
        raise ElaborationException(
            f"Expected PDF {output_pdf_path} not found.")

    print(f"PDF successfully created at {output_pdf_path}")

    return output_pdf_path


def get_folder_prs_images_from_pptx(usermail, project, filename, temp_folder):
    zip_byte_data_io = download_file_from_s3(
        usermail, project, filename)  # Assuming this downloads a zip file

    # Unzip to get the pptx file
    pptx_file_path, temp_folder = unzip_file_to_temp(
        zip_byte_data_io, usermail, project)

    # You've unzipped the file, so there's no need to save the buffer again
    prs = Presentation(pptx_file_path)  # Load the presentation from the path

    # Convert the PPTX to PDF and then to images
    pdf_path = pptx_to_pdf(pptx_file_path).replace('.pptx.pdf', '.pdf')
    print(f"pdf_path: {pdf_path}")
    images = pdf_to_images(pdf_path)
    print(f"images: {images}")

    return prs, images


def generate_audio(index, folder, prefix, text, voice_name, base64):
    filename = os.path.join(folder, f'{prefix}_{index}.mp3')
    generate_tts(text, voice_name, filename)
    audio = AudioSegment.from_file(filename)
    if (base64):
        return audiosegment_to_base64(audio)
    return audio, filename


def combine_audios_and_generate_file(index, folder, audios, base64):
    audios.pop()
    combined_audio = audios[0]
    for audio in audios[1:]:
        combined_audio += audio
    combined_filename = os.path.join(folder, f'slide_{index}.mp3')
    combined_audio.export(combined_filename, format="mp3", bitrate="320k")
    if (base64):
        return audiosegment_to_base64(combined_audio)
    return combined_audio, combined_filename

# PROCESS PPTX BLOCK
def process_pptx(usermail, project, filename):
    zip_byte_data = download_file_from_s3(usermail, project, filename)
    pptx_file, temp_folder = unzip_file_to_temp(
        zip_byte_data, usermail, project)
    
    if add_tts_to_pptx(pptx_file, temp_folder):
        edited_filename = f"{os.path.splitext(filename)[0]}_edited.zip"
        edited_filepath = zip_pptx(pptx_file, temp_folder, edited_filename)
        print(edited_filepath)
        with open(edited_filepath, 'rb') as file_to_upload:
            upload_file_to_s3(usermail, project,
                              edited_filename, file_to_upload)
        try:
            delete_folder(temp_folder)
        except Exception as e:
            raise ElaborationException(
                f"Exception while deleting temp folder: {str(e)}")
    else:
        raise ElaborationException("PPTX has no notes to elaborate")


def add_tts_to_pptx(pptx_file, temp_folder):
    try:
        prs = Presentation(pptx_file)
        modified = False
        for slide in prs.slides:
            if process_slide(slide, temp_folder):
                modified = True
        if modified:
            prs.save(pptx_file)
        return modified
    except AmazonException as e:
        raise AmazonException(e)
    except UserParameterException as e:
        raise UserParameterException(e)
    except ElaborationException as e:
        raise ElaborationException(e)
    except Exception as e:
        raise ElaborationException(f"Exception adding tts to pptx: {str(e)}")


def process_slide(slide, temp_folder):
    try:
        notes_text, modified = check_slide_have_notes(slide.notes_slide)
        if not modified:
            return modified
        parsed_ssml = check_correct_validate_parse_text(notes_text)
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise ElaborationException(
            f"Presentation Elaboration: Exception during SSML correction/validation/parsing: {str(e)}")

    try:
        if len(parsed_ssml) == 1:
            unique_id = uuid.uuid4()
            for voice_name, text in parsed_ssml:
                audio, filename = generate_audio(
                    unique_id, temp_folder, "slide", text, voice_name, False)
                slide.shapes.turbo_add_enabled = True
                add_audio_to_slide(slide, filename)
                slide.shapes.turbo_add_enabled = False
        else:
            audios = []
            for voice_name, text in parsed_ssml:
                unique_id = uuid.uuid4()
                audio, _ = generate_audio(unique_id, temp_folder,
                                          "multi_voice", text, voice_name, False)
                audios.append(audio)
                audios.append(half_sec_silence)
            combined_audio, combined_filename = combine_audios_and_generate_file(
                "combined", temp_folder, audios, False)
            slide.shapes.turbo_add_enabled = True
            add_audio_to_slide(slide, combined_filename)
            slide.shapes.turbo_add_enabled = False
        return modified
    except AmazonException as e:
        raise AmazonException(e)
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise ElaborationException(
            f"Exception while processing a slide: {str(e)}")

# TTS TEXT PREVIEW BLOCK


def process_preview(text):
    unique_id = uuid.uuid4()
    try:
        parsed_ssml = check_correct_validate_parse_text(text)
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise ElaborationException(
            f"Audio Preview: Exception during SSML correction/validation/parsing: {str(e)}")
    try:
        temp_folder = create_folder(f"{unique_id}_temp")
        if len(parsed_ssml) == 1:
            for voice_name, text in parsed_ssml:
                audio, _ = generate_audio(
                    voice_name, temp_folder, "slide", text, voice_name, False)
                audio_buffer = audiosegment_to_stream(audio)
                return audio_buffer
        else:
            audios = []
            for voice_name, text in parsed_ssml:
                audio, _ = generate_audio(voice_name, temp_folder,
                                          "multi_voice", text, voice_name, False)
                audios.append(audio)
                audios.append(half_sec_silence)
            combined_audio, _ = combine_audios_and_generate_file(
                "combined", temp_folder, audios, False)
            audio_buffer = audiosegment_to_stream(combined_audio)

            return audio_buffer
    except AmazonException as e:
        raise AmazonException(e)
    except ElaborationException as e:
        raise ElaborationException(e)
    except UserParameterException as e:
        raise UserParameterException(e)
    except Exception as e:
        raise ElaborationException(
            f"Exception while processing text: {str(e)}")
    finally:
        delete_folder(temp_folder)

# PPTX SPLIT BLOCK
def get_slide_audio_preview(index, slide, folder):
    notes_text, have_notes = check_slide_have_notes(slide.notes_slide)
    if not have_notes:
        return None
    try:
        parsed_ssml = check_correct_validate_parse_text(notes_text)
        audio_mp3 = None
        if len(parsed_ssml) == 1:
            for voice_name, text in parsed_ssml:
                audio_mp3, _ = generate_audio(
                    index, folder, "slide", text, voice_name, False)
        else:
            audios = []
            for j, (voice_name, text) in enumerate(parsed_ssml):
                audio, _ = generate_audio(
                    j, folder, "multi_voice", text, voice_name, False)
                audios.append(audio)
                audios.append(half_sec_silence)
            audio_mp3, _ = combine_audios_and_generate_file(
                index, folder, audios, False)
        return audio_mp3
    except AmazonException as e:
        print(f"Amazon Exception")
        return None
    except Exception as e:
        print(f"Exception while saving audio to file: {str(e)}")
        return None


def process_pptx_split(usermail, project, filename):
    try:
        temp_folder = create_folder(f"{usermail}_{project}_temp")
        prs, images = get_folder_prs_images_from_pptx(
            usermail, project, filename, temp_folder)

        stream = io.BytesIO()
        with ZipFile(stream, "w") as zf:
            for i, (slide, image) in enumerate(zip(prs.slides, images)):
                print(f"Processing slide {i}")
                audio_segment = get_slide_audio_preview(i, slide, temp_folder)
                img_buffer = io.BytesIO()
                print(img_buffer)
                image.save(img_buffer, 'PNG')
                print(f"Saving slide {i} to zip")
                img_buffer.seek(0)
                zf.writestr(f"slide_{i}.png", img_buffer.read())
                print(f"Saving audio {i} tow zip")
                if (audio_segment):
                    audio_buffer = audiosegment_to_stream(audio_segment)
                    zf.writestr(f"audio_{i}.mp3", audio_buffer.read())
        stream.seek(0)
        return stream
    except OSError as e:
        raise ElaborationException(
            f"Could not process file due to OS path issue: {e.filename}")
    except AmazonException as e:
        raise AmazonException(e)
    except UserParameterException as e:
        raise UserParameterException(e)
    except ElaborationException as e:
        raise ElaborationException(e)
    except Exception as e:
        # traceback.print_exc()
        raise ElaborationException(
            f"Exception while processing slides splitted: {str(e)}")
    finally:
        delete_folder(temp_folder)
